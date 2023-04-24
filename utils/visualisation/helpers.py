""" General helper functions common across visualisations.
"""

import matplotlib
import matplotlib.pyplot as plt

import numpy
import collections

import beautifulcells.utils.inputHelpers as input
import beautifulcells.utils.calculations as calcs

def dealWithPlot(savePlot, showPlot, closePlot, folder, plotName, dpi,
				 tightLayout=True):
	""" Deals with the current matplotlib.pyplot.
	"""

	if tightLayout:
		plt.tight_layout()

	if savePlot:
		plt.savefig(folder+plotName, dpi=dpi,
					format=plotName.split('.')[-1])

	if showPlot:
		plt.show()

	if closePlot:
		plt.close()

def getFigAxes(fig, axes, figsize, **kwargs):
	if type(fig) == type(None):
		fig = plt.figure(figsize=figsize, **kwargs)
	if type(axes) == type(None):
		axes = fig.add_subplot(111)  # 1 row, 1 col, start at cell 1.

	return fig, axes

def getColors(labels, labelSet=None, colorMap='tab20', rgb=False):
	""" Gets an OrderedDict of colors; the order indicates the frequency of \
	labels from largest to smallest.

	Args:
		labels (numpy.array<str>): Indicates a set of labels for observations.

		labelSet (list-like<str>): Indicates the set of labels in labels. \
									If None, calculated based on labels.

		colorMap (str): A matplotlib colormap.

		rgb (bool): If True, colors indicated by rgb value, if false hexcode.

	Returns:
		dict<str, tuple or str>: An ordered dict indicating the labels which \
					occur most to least frequently and the associated colors.
	"""
	# Determining the set of labels #
	labelSet = input.returnDefaultIfNone(labelSet,
										 calcs.getOrderedLabelSet(labels))

	# Initialising the ordered dict #
	cellTypeColors = {}

	# Ordering the cells according to their frequency and obtaining colors #
	nLabels = len(labelSet)
	cmap = plt.cm.get_cmap(colorMap, nLabels)
	rgbs = [cmap(i)[:3] for i in range(nLabels)]
	#rgbs = list(numpy.array(rgbs)[order]) # Make sure color order is the same.

	# Populating the color dictionary with rgb values or hexcodes #
	for i in range(len(labelSet)):
		cellType = labelSet[i]
		rgbi = rgbs[i]
		if not rgb:
			cellTypeColors[cellType] = matplotlib.colors.rgb2hex(rgbi)
		else:
			cellTypeColors[cellType] = rgbi

	return cellTypeColors

"""
# Example of how to make powerpoint from png images #
import pptx
from pptx.util import Inches

prs = pptx.Presentation()
width, height = 16, 9
prs.slide_width = Inches(width)
prs.slide_height = Inches(height)
blank_slide_layout = prs.slide_layouts[6]  # choosing a slide layout

# Adding the cluster results #
slide = prs.slides.add_slide(blank_slide_layout)
imgs = ['foxg1_clusters-wt.png', 'foxg1_clusters-ko.png']
rows = [Inches(.1)]
cols = [Inches(.1), Inches(width/2)]
for i, img in enumerate(imgs):
    slide.shapes.add_picture(out_plots+img, cols[i], rows[0])

# Adding the gene expression #
genes_per_slide = [['Hcrt', 'Foxg1'], ['Ddc', 'Th', 'Slc6a3'], ['Tac2', 'Pax6']]
for i, genes in enumerate(genes_per_slide):
    slide = prs.slides.add_slide(blank_slide_layout)

    rows = [Inches(.1), Inches(height/2)]
    cols = [Inches(.1)]+[Inches((width/len(genes))*j)
                         for j in range(1,len(genes))]

    for j, gene in enumerate(genes):
        img_wt = out_plots+'gene_plots/'+f'foxg1_{gene}-wt_umap.png'
        img_ko = out_plots+'gene_plots/'+f'foxg1_{gene}-ko_umap.png'
        slide.shapes.add_picture(img_wt, cols[j], rows[0])
        slide.shapes.add_picture(img_ko, cols[j], rows[1])

prs.save(out_dir+"foxg1_np-day_analysis.pptx") # saving file
"""

